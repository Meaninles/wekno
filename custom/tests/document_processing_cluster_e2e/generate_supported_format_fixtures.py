from __future__ import annotations

import argparse
import csv
import json
import shutil
import subprocess
from email import policy as email_policy
from email.message import EmailMessage
from pathlib import Path

from docx import Document
from ebooklib import epub
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt


SUPPORTED_EXTENSIONS = (
    "pdf",
    "txt",
    "text",
    "docx",
    "doc",
    "epub",
    "mhtml",
    "md",
    "markdown",
    "png",
    "jpg",
    "jpeg",
    "gif",
    "webp",
    "bmp",
    "tiff",
    "csv",
    "xlsx",
    "xls",
    "pptx",
    "ppt",
    "json",
    "mp3",
    "wav",
    "m4a",
    "flac",
    "ogg",
)


def marker(extension: str) -> str:
    return f"WKN-FORMAT-{extension.upper()}-7319"


def prose(extension: str) -> str:
    return (
        f"{marker(extension)}。该文件用于验证 {extension.upper()} 文档的解析、切分、"
        "向量化、摘要、问题生成、知识图谱和 Wiki 衍生流程。"
        "审批负责人必须在三个工作日内完成复核，未经授权不得跳过安全检查。"
    )


def run(command: list[str]) -> None:
    completed = subprocess.run(
        command,
        check=False,
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
        text=True,
    )
    if completed.returncode != 0:
        raise RuntimeError(f"command failed ({completed.returncode}): {command}\n{completed.stdout}")


def write_simple_formats(output: Path) -> list[Path]:
    paths: list[Path] = []
    for extension in ("txt", "text", "md", "markdown"):
        path = output / f"format-{extension}.{extension}"
        prefix = "# 支持格式验收\n\n" if extension in {"md", "markdown"} else "支持格式验收\n"
        path.write_text(prefix + prose(extension) + "\n", encoding="utf-8")
        paths.append(path)

    csv_path = output / "format-csv.csv"
    with csv_path.open("w", encoding="utf-8-sig", newline="") as stream:
        writer = csv.writer(stream)
        writer.writerow(["标识", "责任部门", "完成时限", "控制要求"])
        writer.writerow([marker("csv"), "数字化管理部门", "三个工作日", "未经授权不得跳过安全检查"])
    paths.append(csv_path)

    json_path = output / "format-json.json"
    json_path.write_text(
        json.dumps(
            {
                "marker": marker("json"),
                "responsible_department": "数字化管理部门",
                "deadline": "三个工作日",
                "control": "未经授权不得跳过安全检查",
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    paths.append(json_path)
    return paths


def write_office_formats(output: Path) -> list[Path]:
    paths: list[Path] = []

    def save_document(extension: str, path: Path) -> None:
        document = Document()
        document.add_heading("支持格式验收", 0)
        document.add_paragraph(prose(extension))
        document.add_heading("职责与时限", level=1)
        document.add_paragraph("数字化管理部门负责审批，负责人必须在三个工作日内完成复核。")
        document.save(path)

    def save_workbook(extension: str, path: Path) -> None:
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "验收规则"
        sheet.append(["标识", "责任部门", "完成时限", "控制要求"])
        sheet.append([marker(extension), "数字化管理部门", "三个工作日", "未经授权不得跳过安全检查"])
        workbook.save(path)

    def save_presentation(extension: str, path: Path) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "支持格式验收"
        textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.0), Inches(3.0))
        frame = textbox.text_frame
        frame.text = prose(extension)
        frame.paragraphs[0].font.size = Pt(24)
        presentation.save(path)

    docx_path = output / "format-docx.docx"
    save_document("docx", docx_path)
    paths.append(docx_path)

    xlsx_path = output / "format-xlsx.xlsx"
    save_workbook("xlsx", xlsx_path)
    paths.append(xlsx_path)

    pptx_path = output / "format-pptx.pptx"
    save_presentation("pptx", pptx_path)
    paths.append(pptx_path)

    conversion_root = output / "office-conversion"
    conversion_root.mkdir(exist_ok=True)
    doc_source = conversion_root / "format-doc.docx"
    pdf_source = conversion_root / "format-pdf.docx"
    xls_source = conversion_root / "format-xls.xlsx"
    ppt_source = conversion_root / "format-ppt.pptx"
    save_document("doc", doc_source)
    save_document("pdf", pdf_source)
    save_workbook("xls", xls_source)
    save_presentation("ppt", ppt_source)
    run(["libreoffice", "--headless", "--convert-to", "doc", "--outdir", str(conversion_root), str(doc_source)])
    run(["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", str(conversion_root), str(pdf_source)])
    run(["libreoffice", "--headless", "--convert-to", "xls", "--outdir", str(conversion_root), str(xls_source)])
    run(["libreoffice", "--headless", "--convert-to", "ppt", "--outdir", str(conversion_root), str(ppt_source)])
    conversion_names = {
        "format-doc.doc": output / "format-doc.doc",
        "format-pdf.pdf": output / "format-pdf.pdf",
        "format-xls.xls": output / "format-xls.xls",
        "format-ppt.ppt": output / "format-ppt.ppt",
    }
    for converted_name, target in conversion_names.items():
        source = conversion_root / converted_name
        if not source.exists():
            raise RuntimeError(f"LibreOffice did not create {source}")
        shutil.move(str(source), target)
        paths.append(target)
    shutil.rmtree(conversion_root)
    return paths


def write_web_formats(output: Path) -> list[Path]:
    epub_path = output / "format-epub.epub"
    book = epub.EpubBook()
    book.set_identifier("weknora-supported-format-epub")
    book.set_title("支持格式验收")
    book.set_language("zh-CN")
    chapter = epub.EpubHtml(title="验收规则", file_name="rules.xhtml", lang="zh-CN")
    chapter.content = (
        "<html><body><h1>支持格式验收</h1>"
        f"<p>{prose('epub')}</p></body></html>"
    )
    book.add_item(chapter)
    book.toc = (chapter,)
    book.spine = ["nav", chapter]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    epub.write_epub(str(epub_path), book)

    mhtml_path = output / "format-mhtml.mhtml"
    message = EmailMessage()
    message["Subject"] = "支持格式验收"
    message["MIME-Version"] = "1.0"
    message.set_content(
        f"<html><body><h1>支持格式验收</h1><p>{prose('mhtml')}</p></body></html>",
        subtype="html",
        charset="utf-8",
    )
    mhtml_path.write_bytes(message.as_bytes(policy=email_policy.default))
    return [epub_path, mhtml_path]


def load_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = (
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf",
    )
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size)
    return ImageFont.load_default()


def write_image_formats(output: Path) -> list[Path]:
    paths: list[Path] = []
    title_font = load_font(54)
    body_font = load_font(34)
    for extension in ("png", "jpg", "jpeg", "gif", "webp", "bmp", "tiff"):
        image = Image.new("RGB", (1600, 900), "white")
        draw = ImageDraw.Draw(image)
        draw.rectangle((50, 50, 1550, 850), outline="#2563eb", width=8)
        draw.text((110, 120), "WEKNORA FORMAT ACCEPTANCE", fill="black", font=title_font)
        draw.text((110, 260), marker(extension), fill="#b91c1c", font=title_font)
        draw.text((110, 410), "APPROVAL OWNER: DIGITAL MANAGEMENT", fill="black", font=body_font)
        draw.text((110, 500), "DEADLINE: THREE WORKDAYS", fill="black", font=body_font)
        draw.text((110, 590), "SECURITY CHECK MUST NOT BE SKIPPED", fill="black", font=body_font)
        path = output / f"format-{extension}.{extension}"
        format_name = "JPEG" if extension in {"jpg", "jpeg"} else extension.upper()
        save_options: dict[str, object] = {}
        if format_name == "JPEG":
            save_options["quality"] = 95
        if format_name == "GIF":
            image = image.quantize(colors=256)
        image.save(path, format=format_name, **save_options)
        paths.append(path)
    return paths


def write_audio_formats(output: Path, source_wav: Path) -> list[Path]:
    if not source_wav.is_file():
        raise RuntimeError(f"audio source does not exist: {source_wav}")
    paths: list[Path] = []
    wav_target = output / "format-wav.wav"
    shutil.copyfile(source_wav, wav_target)
    paths.append(wav_target)
    codecs = {
        "mp3": ["-codec:a", "libmp3lame", "-q:a", "3"],
        "m4a": ["-codec:a", "aac", "-b:a", "128k"],
        "flac": ["-codec:a", "flac"],
        "ogg": ["-codec:a", "libvorbis", "-q:a", "5"],
    }
    for extension, options in codecs.items():
        target = output / f"format-{extension}.{extension}"
        run(["ffmpeg", "-hide_banner", "-loglevel", "error", "-y", "-i", str(source_wav), *options, str(target)])
        paths.append(target)
    return paths


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("output", type=Path)
    parser.add_argument("--audio-source", required=True, type=Path)
    args = parser.parse_args()
    output = args.output.resolve()
    if output.exists():
        shutil.rmtree(output)
    output.mkdir(parents=True)

    paths = [
        *write_simple_formats(output),
        *write_office_formats(output),
        *write_web_formats(output),
        *write_image_formats(output),
        *write_audio_formats(output, args.audio_source.resolve()),
    ]
    by_extension = {path.suffix.lstrip(".").lower(): path for path in paths}
    missing = sorted(set(SUPPORTED_EXTENSIONS) - set(by_extension))
    extra = sorted(set(by_extension) - set(SUPPORTED_EXTENSIONS))
    if missing or extra or len(paths) != len(SUPPORTED_EXTENSIONS):
        raise RuntimeError(
            f"fixture extension mismatch: missing={missing}, extra={extra}, paths={len(paths)}"
        )
    manifest = {
        "schema_version": 1,
        "files": [
            {
                "extension": extension,
                "filename": by_extension[extension].name,
                "marker": marker(extension),
                "expected_chunk_text": (
                    ""
                    if extension in {"mp3", "wav", "m4a", "flac", "ogg"}
                    else marker(extension)
                ),
                "expected_derivatives": (
                    ["multimodal"]
                    if extension in {"png", "jpg", "jpeg", "gif", "webp", "bmp", "tiff"}
                    else ["table"]
                    if extension in {"csv", "xlsx", "xls"}
                    else []
                ),
                "size_bytes": by_extension[extension].stat().st_size,
            }
            for extension in SUPPORTED_EXTENSIONS
        ],
    }
    (output / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    print(json.dumps(manifest, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
